import io
from pptx import Presentation
from typing import List
from api.schemas import DocumentChunk, DocumentMetadata
from processing.cleaner import clean_text
from processing.chunker import chunk_text

def load_pptx(file_bytes: bytes, file_name: str, document_id: str) -> List[DocumentChunk]:
    chunks = []
    
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
                    
            text = "\n".join(slide_text)
            cleaned_text = clean_text(text)
            
            if cleaned_text:
                metadata = DocumentMetadata(
                    file_name=file_name,
                    file_type="pptx",
                    slide_number=i + 1,
                    page_number=i + 1, # use slide number as page number for generic UI mapping
                    document_id=document_id
                )
                slide_chunks = chunk_text(cleaned_text, metadata)
                chunks.extend(slide_chunks)
                
    except Exception as e:
        print(f"Failed to process PPTX {file_name}: {e}")
        
    return chunks
